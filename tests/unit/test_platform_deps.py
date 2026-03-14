"""Tests for S9b.2 — Platform Dependency Declaration.

Verifies that all new platform dependencies are importable and functional.
"""

from __future__ import annotations


def test_anthropic_sdk_importable():
    """Anthropic SDK for P22 code generation."""
    import anthropic

    assert hasattr(anthropic, "Anthropic")
    assert hasattr(anthropic, "AsyncAnthropic")


def test_jose_jwt_importable():
    """python-jose for JWT token handling (P14 auth)."""
    from jose import jwt

    assert callable(jwt.encode)
    assert callable(jwt.decode)


def test_passlib_bcrypt_importable():
    """passlib for password hashing (P14 auth)."""
    from passlib.hash import bcrypt

    assert callable(bcrypt.hash)
    assert callable(bcrypt.verify)


def test_websockets_importable():
    """websockets for real-time comments (P14)."""
    import websockets

    assert hasattr(websockets, "connect")


def test_edge_tts_importable():
    """edge-tts for text-to-speech (P23 audio)."""
    import edge_tts

    assert hasattr(edge_tts, "Communicate")


def test_python_pptx_importable():
    """python-pptx for slide generation (P19)."""
    from pptx import Presentation

    assert callable(Presentation)


def test_bcrypt_hash_roundtrip():
    """Verify bcrypt hashing works end-to-end."""
    import bcrypt

    password = b"test-password-123"
    hashed = bcrypt.hashpw(password, bcrypt.gensalt())
    assert bcrypt.checkpw(password, hashed)
    assert not bcrypt.checkpw(b"wrong-password", hashed)


def test_jose_jwt_roundtrip():
    """Verify JWT encode/decode works end-to-end."""
    from jose import jwt

    secret = "test-secret-key"
    payload = {"sub": "user@example.com", "role": "researcher"}
    token = jwt.encode(payload, secret, algorithm="HS256")
    decoded = jwt.decode(token, secret, algorithms=["HS256"])
    assert decoded["sub"] == "user@example.com"
    assert decoded["role"] == "researcher"


def test_python_pptx_create_presentation():
    """Verify PPTX creation works."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    assert slide is not None
    assert len(prs.slides) == 1
